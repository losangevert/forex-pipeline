#!/usr/bin/env python3
"""
Génère le PowerPoint de soutenance pour le projet Forex Pipeline.
3 orateurs, 10 min de présentation + 5 min de démo = 15 min.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Couleurs (thème clair)
BG = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x15, 0x65, 0xC0)
DARK = RGBColor(0x1a, 0x1a, 0x2e)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x5C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHTGRAY = RGBColor(0x88, 0x88, 0x88)
SECTION = RGBColor(0xF5, 0xF5, 0xF5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_bullet_slide(slide, items, left=0.8, top=1.8, width=11.5, size=20, color=DARK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return txBox

def speaker_tag(slide, name, left=0.8, top=0.3):
    tag = add_shape(slide, Inches(left), Inches(top), Inches(2.5), Inches(0.4), ORANGE)
    tag.text_frame.paragraphs[0].text = f"  🎤 {name}"
    tag.text_frame.paragraphs[0].font.size = Pt(14)
    tag.text_frame.paragraphs[0].font.color.rgb = WHITE
    tag.text_frame.paragraphs[0].font.bold = True
    tag.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

def slide_number(slide, num, total=10):
    add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                 f"{num}/{total}", size=12, color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 1 — Titre
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), BLUE)
add_shape(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), BLUE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "Suivi des taux de change multi-devises", size=40, color=BLUE, bold=True)
add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.6),
             "Plateforme d'orchestration Airflow • Frankfurter API • Metabase", size=22, color=BLUE)

add_shape(slide, Inches(1), Inches(3.8), Inches(3), Inches(0.03), ORANGE)

add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
             "Projet final — Cours IA", size=18, color=GRAY)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
             "👤 Lucas • Orateur 2 • Orateur 3", size=18, color=GRAY)
add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
             "Juin 2026", size=16, color=GRAY)

# ============================================================
# SLIDE 2 — Contexte & Objectifs
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
             "Contexte & Objectifs", size=32, color=BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(3), Inches(0.04), BLUE)
speaker_tag(slide, "Orateur 1 — 2 min")
slide_number(slide, 2)

add_bullet_slide(slide, [
    "🎯 Objectif :",
    "      Plateforme automatisée d'ingestion et d'analyse",
    "      des taux de change multi-devises",
    "",
    "📊 API utilisée : Frankfurter (taux ECB quotidiens)",
    "",
    "🔄 Pipeline complet :",
    "      Extraction → Stockage → Transformation → QC → Analyse",
    "",
    "📈 Visualisation : Metabase pour les KPIs métier",
    "",
    "⚙️ Stack technique : Airflow 2.10.5 / PostgreSQL 16 / Metabase",
    "      Docker Compose • VPS OVH • GitHub",
], size=20, color=DARK)

# ============================================================
# SLIDE 3 — Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
             "Architecture technique", size=32, color=BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(3), Inches(0.04), BLUE)
speaker_tag(slide, "Orateur 1 — 1 min")
slide_number(slide, 3)

# Draw architecture boxes
boxes = [
    (1.0, 2.0, 2.5, 1.2, "Frankfurter API", "https://api.frankfurter.app", BLUE),
    (4.5, 2.0, 2.5, 1.2, "Airflow Scheduler", "DAG @hourly", GREEN),
    (8.0, 2.0, 2.5, 1.2, "PostgreSQL 16", "forex DB", ORANGE),
    (4.5, 4.0, 2.5, 1.2, "Airflow Webserver", "UI • port 8080", GREEN),
    (8.0, 4.0, 2.5, 1.2, "Metabase", "Dashboards • port 3000", ORANGE),
]

for l, t, w, h, title, sub, color in boxes:
    shape = add_shape(slide, Inches(l), Inches(t), Inches(w), Inches(h), color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = sub
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
    p2.alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.8),
             "Docker Compose sur VPS OVH (Debian trixie) — 7.6 Go RAM, 74 Go disque",
             size=14, color=GRAY)

# ============================================================
# SLIDE 4 — DAG Airflow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
             "Pipeline Airflow (DAG)", size=32, color=BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(3), Inches(0.04), BLUE)
speaker_tag(slide, "Orateur 2 — 2 min")
slide_number(slide, 4)

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
             "forex_exchange_rate_pipeline — 7 tâches (TaskFlow API @dag/@task)",
             size=18, color=GRAY)

# DAG steps as boxes
steps = [
    (0.8, 2.5, 1.5, 0.8, "1. extract_raw", "API Frankfurter"),
    (2.7, 2.5, 1.5, 0.8, "2. store_raw", "JSONB → raw_rates"),
    (4.6, 2.5, 1.8, 0.8, "3. transform_validate", "QC 5 dimensions"),
    (6.8, 2.5, 1.5, 0.8, "4. load_valid", "UPSERT → exchange_rates"),
    (8.7, 2.5, 1.5, 0.8, "5. load_graveyard", "Lignes invalides"),
    (10.6, 2.5, 1.8, 0.8, "6. detect_anomalies", "Seuil 2%"),
]

for i, (l, t, w, h, title, sub) in enumerate(steps):
    color = BLUE if i < 3 else GREEN
    shape = add_shape(slide, Inches(l), Inches(t), Inches(w), Inches(h), color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = sub
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
    p2.alignment = PP_ALIGN.CENTER

# 7th step: log_pipeline
shape = add_shape(slide, Inches(5.5), Inches(4.0), Inches(2.0), Inches(0.8), ORANGE)
tf = shape.text_frame
p = tf.paragraphs[0]
p.text = "7. log_pipeline"
p.font.size = Pt(14)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Statut complet"
p2.font.size = Pt(10)
p2.font.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
p2.alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(1.5),
             "⚙️ Config externalisée : forex_currencies, forex_alert_threshold, forex_api_base\n"
             "🔄 Idempotent : ON CONFLICT DO NOTHING\n"
             "🛡️ Retries + timeouts : extract_raw (x2), load_valid (x2), log_pipeline (x2)",
             size=16, color=GRAY)

# ============================================================
# SLIDE 5 — Contrôle qualité
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
             "Contrôle qualité des données", size=32, color=BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(3), Inches(0.04), BLUE)
speaker_tag(slide, "Orateur 2 — 1 min 30")
slide_number(slide, 5)

qc_items = [
    "✅  Complétude   — Vérifie que chaque paire attendue est présente dans la réponse API",
    "",
    "✅  Cohérence    — Les taux sont des nombres > 0 (sinon → cimetière)",
    "",
    "✅  Fraîcheur    — La date ne dépasse pas le seuil (Variable forex_freshness_hours)",
    "",
    "✅  Unicité      — Contrainte UNIQUE (currency_pair, rate_date) + ON CONFLICT",
    "",
    "✅  Structure    — Vérification de la clé 'rates' et du format JSON",
    "",
    "🗑️  Lignes rejetées → data_quality_graveyard (cimetière tracé)",
    "",
    "📊  Pipeline log : 1 entrée par exécution (status, reçues, valides, rejetées, insérées)",
]
add_bullet_slide(slide, qc_items, size=17, color=DARK)

# ============================================================
# SLIDE 6 — Alertes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
             "Détection d'anomalies (Alertes)", size=32, color=BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(3), Inches(0.04), BLUE)
speaker_tag(slide, "Orateur 2 — 1 min")
slide_number(slide, 6)

add_bullet_slide(slide, [
    "🔍  Comparaison entre le taux actuel et le taux précédent en base",
    "",
    "📐  Calcul : |taux_n - taux_n-1| / taux_n-1 × 100",
    "",
    "⚡  Seuil déclencheur : 2 % (configurable via Variable)",
    "",
    "📝  Justification du seuil à 2 % :",
    "       • Les paires EUR/USD, EUR/GBP varient de 0,1 % à 1 % / jour",
    "       • Un seuil à 2 % capte les mouvements anormaux sans faux positifs",
    "       • Exemples : crise, annonce macroéconomique, décision de banque centrale",
    "",
    "💾  Écriture dans la table rate_alerts (devise, avant, après, écart)",
], size=18, color=DARK)

# ============================================================
# SLIDE 7 — KPIs Metabase
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
             "KPIs & Dashboards Metabase", size=32, color=BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(3), Inches(0.04), BLUE)
speaker_tag(slide, "Orateur 3 — 1 min 30")
slide_number(slide, 7)

add_bullet_slide(slide, [
    "📈  Dashboard « Taux de change — Suivi »",
    "       • Courbes individuelles EUR/USD, EUR/GBP, EUR/JPY, EUR/CHF, EUR/AUD",
    "       • Vue superposée pour comparer toutes les paires",
    "       • Bar chart des top variations hebdomadaires",
    "",
    "📋  Dashboard « Pipeline — Monitoring »",
    "       • Tableau des 10 derniers logs d'exécution",
    "       • Tableau des alertes de variation déclenchées",
    "",
    "🗄️  Vues SQL disponibles :",
    "       • v_last_30d_trend : tendance + variation quotidienne",
    "       • v_top_weekly_variations : top 20 variations absolues",
    "",
    "🎯  Questions libres possibles : requêtes SQL personnalisées",
], size=18, color=DARK)

# ============================================================
# SLIDE 8 — Robustesse
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
             "Choix de robustesse", size=32, color=BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(3), Inches(0.04), BLUE)
speaker_tag(slide, "Orateur 3 — 1 min")
slide_number(slide, 8)

add_bullet_slide(slide, [
    "🔄  Retries adaptés par tâche (1 à 2 selon le risque)",
    "",
    "⏱️  Timeouts : API (60s), transformations (120s), logs (30s)",
    "",
    "🛡️  Chemins nominal ET d'échec :",
    "       • API injoignable → retry → fail tracé",
    "       • Ligne invalide à l'insertion → cimetière",
    "",
    "📦  Idempotence : ON CONFLICT DO NOTHING",
    "",
    "🔧  Configuration externalisée : Variables Airflow",
    "       (devises, seuils, URL API — sans redéploiement)",
    "",
    "📝  Log complet de chaque exécution (pipeline_log)",
], size=18, color=DARK)

# ============================================================
# SLIDE 9 — Démo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
             "Démonstration", size=32, color=BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(3), Inches(0.04), ORANGE)
speaker_tag(slide, "Équipe — 5 min")
slide_number(slide, 9)

add_bullet_slide(slide, [
    "1.  Airflow UI — Graph View du DAG",
    "       Affichage des 7 tâches et de leurs dépendances",
    "",
    "2.  Exécution du pipeline",
    "       Trigger manuel + suivi en temps réel",
    "",
    "3.  Vérification des données PostgreSQL",
    "       Tables : exchange_rates, pipeline_log, rate_alerts",
    "",
    "4.  Dashboards Metabase",
    "       Graphiques d'évolution + monitoring pipeline",
], size=20, color=DARK)

# ============================================================
# SLIDE 10 — Conclusion
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), BLUE)
add_shape(slide, Inches(0), Inches(7.35), Inches(13.333), Inches(0.15), BLUE)

add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
             "Conclusion", size=32, color=BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(3), Inches(0.04), BLUE)
slide_number(slide, 10)

add_bullet_slide(slide, [
    "✅  Pipeline complet et fonctionnel",
    "",
    "✅  Contrôle qualité automatisé (5 dimensions)",
    "",
    "✅  Alertes configurables",
    "",
    "✅  KPIs visuels dans Metabase",
    "",
    "✅  Déploiement Docker sur VPS OVH",
    "",
    "📦  Code source : https://github.com/losangevert/forex-pipeline",
], size=20, color=DARK)

# Questions slide
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8),
             "🙋 Questions ?", size=28, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# Save
output_path = "/home/node/.openclaw/workspace/airflow-forex-pipeline/presentation_forex_pipeline.pptx"
prs.save(output_path)
print(f"✅ Présentation sauvegardée : {output_path}")
print(f"   {len(prs.slides)} slides — 10 min présentation + 5 min démo")
